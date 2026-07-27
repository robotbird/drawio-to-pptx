#!/usr/bin/env python3
"""
drawio_to_pptx.py - Convert a .drawio diagram into a NATIVE, EDITABLE PowerPoint (.pptx).

Every drawio vertex becomes an editable PPT autoshape / text box; every drawio edge
becomes an editable connector (straight line or open polyline) with arrowheads.
The output is NOT a flat image - all shapes can be moved, recolored, re-texted and
re-routed in PowerPoint / Keynote / WPS.

Usage:
    drawio_to_pptx.py input.drawio -o output.pptx
    drawio_to_pptx.py input.drawio -o out.pptx --size 16:9 --font-scale 1.2
    drawio_to_pptx.py input.drawio -o out.pptx --landscape

Requires: python-pptx, lxml   (pip install python-pptx lxml)
"""
import argparse
import re
import sys
import zlib
import base64

from lxml import etree

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

EMU_PER_IN = 914400
EMU_PER_PT = 12700

# ---------------------------------------------------------------------------
# drawio XML loading (handles <mxfile>, compressed payloads, and bare models)
# ---------------------------------------------------------------------------

def load_diagrams(path):
    """Return list of (name, <mxGraphModel> element)."""
    parser = etree.XMLParser(recover=True, huge_tree=True)
    tree = etree.parse(str(path), parser)
    root = tree.getroot()
    tag = etree.QName(root).localname

    if tag == "mxGraphModel":
        return [(root.get("name") or "Diagram", root)]

    if tag == "mxfile":
        out = []
        for diag in root.iterfind("*"):
            if etree.QName(diag).localname != "diagram":
                continue
            name = diag.get("name", "Diagram")
            model = None
            inner = diag.find("mxGraphModel")
            if inner is not None:
                model = inner
            else:
                # compressed: base64 text child -> raw deflate -> xml
                text = (diag.text or "").strip()
                if text:
                    try:
                        raw = base64.b64decode(text)
                        xml = zlib.decompress(raw, -15)
                        model = etree.fromstring(xml, etree.XMLParser(recover=True, huge_tree=True))
                    except Exception:
                        model = None
            if model is not None:
                out.append((name, model))
        if out:
            return out

    # last resort: any mxGraphModel anywhere
    found = root.find(".//mxGraphModel")
    if found is not None:
        return [("Diagram", found)]
    raise ValueError("No <mxGraphModel> found in the drawio file.")


# ---------------------------------------------------------------------------
# Style helpers
# ---------------------------------------------------------------------------

def parse_style(style_str):
    d = {}
    if not style_str:
        return d
    for part in style_str.split(";"):
        part = part.strip()
        if not part:
            continue
        if "=" in part:
            k, v = part.split("=", 1)
            d[k.strip()] = v.strip()
        else:
            d[part] = "1"  # bare flag, e.g. "ellipse", "rhombus"
    return d


def hexcolor(h, default=None):
    if not h or h == "none":
        return default
    h = h.lstrip("#").strip()
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return default
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return default


def _mso(name, fallback=MSO_SHAPE.RECTANGLE):
    """Look up an MSO_SHAPE member by name, tolerating version differences."""
    try:
        return MSO_SHAPE[name]
    except (KeyError, AttributeError):
        return fallback


def shape_kind(style):
    """Map a drawio style dict to a python-pptx MSO_SHAPE (or None for text-only)."""
    if "ellipse" in style:
        return MSO_SHAPE.OVAL
    if "rhombus" in style:
        return MSO_SHAPE.DIAMOND
    sh = style.get("shape", "")
    mapping = {
        "cylinder": "CAN", "cylinder3": "CAN",
        "triangle": "ISOSCELES_TRIANGLE",
        "hexagon": "HEXAGON",
        "cloud": "CLOUD",
        "parallelogram": "PARALLELOGRAM",
        "trapezoid": "TRAPEZOID",
        "process": "CHEVRON",
    }
    if sh in mapping:
        return _mso(mapping[sh])
    if sh.startswith("mxgraph."):
        return MSO_SHAPE.RECTANGLE  # branded stencils: render as labeled box
    if "text" in style:
        return None
    if "swimlane" in style or style.get("container") == "1":
        return MSO_SHAPE.ROUNDED_RECTANGLE
    if style.get("rounded") == "1":
        return MSO_SHAPE.ROUNDED_RECTANGLE
    return MSO_SHAPE.RECTANGLE


# ---------------------------------------------------------------------------
# HTML-ish label text -> plain lines
# ---------------------------------------------------------------------------

def html_to_lines(html):
    if html is None:
        return []
    s = html
    s = re.sub(r"(?i)<br\s*/?>", "\n", s)
    s = re.sub(r"(?i)</(div|p|li|h[1-6])>", "\n", s)
    s = re.sub(r"(?i)<[^>]+>", "", s)  # strip remaining tags (font, b, i, span...)
    s = s.replace("\r", "")
    return s.split("\n")


# ---------------------------------------------------------------------------
# XML element building helpers (arrowheads, polylines)
# ---------------------------------------------------------------------------

_id_counter = {"n": 1000}


def next_id():
    _id_counter["n"] += 1
    return _id_counter["n"]


def arrow_type(name):
    if not name or name == "none":
        return None
    if name in ("classic", "open", "block", "oval", "diamond", "ERmany", "ERone"):
        return "triangle" if name in ("classic", "block") else name
    return "triangle"


def build_ln_xml(style):
    """Return an <a:ln> element string with fill, width and arrowheads from a drawio edge style."""
    sc = style.get("strokeColor", "#000000")
    rgb = hexcolor(sc, RGBColor(0, 0, 0))
    hexv = "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2])
    width = style.get("strokeWidth", "1")
    try:
        w_emu = int(round(float(width) * EMU_PER_PT))
    except ValueError:
        w_emu = EMU_PER_PT

    head = arrow_type(style.get("startArrow", "none"))
    tail = arrow_type(style.get("endArrow", "classic"))
    ends = ""
    if head:
        ends += '<a:headEnd type="%s" w="med" len="med"/>' % head
    if tail:
        ends += '<a:tailEnd type="%s" w="med" len="med"/>' % tail
    return ('<a:ln w="%d" cap="flat"><a:solidFill><a:srgbClr val="%s"/></a:solidFill>%s</a:ln>'
            % (w_emu, hexv, ends))


def set_connector_arrows(conn, style):
    """Attach arrowheads to a python-pptx connector."""
    spPr = conn._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = spPr.makeelement(qn("a:ln"), {"w": str(EMU_PER_PT)})
        spPr.append(ln)
    # color
    rgb = hexcolor(style.get("strokeColor", "#000000"), RGBColor(0, 0, 0))
    if ln.find(qn("a:noFill")) is not None:
        ln.remove(ln.find(qn("a:noFill")))
    if ln.find(qn("a:solidFill")) is None:
        sf = parse_xml('<a:solidFill %s><a:srgbClr val="%02X%02X%02X"/></a:solidFill>'
                       % (nsdecls("a"), rgb[0], rgb[1], rgb[2]))
        ln.insert(0, sf)
    for tag in ("a:headEnd", "a:tailEnd"):
        e = ln.find(qn(tag))
        if e is not None:
            ln.remove(e)
    head = arrow_type(style.get("startArrow", "none"))
    tail = arrow_type(style.get("endArrow", "classic"))
    if tail:
        ln.append(parse_xml('<a:tailEnd %s type="%s" w="med" len="med"/>' % (nsdecls("a"), tail)))
    if head:
        ln.append(parse_xml('<a:headEnd %s type="%s" w="med" len="med"/>' % (nsdecls("a"), head)))


def add_polyline(slide, points, style):
    """Add an open polyline (for edges with waypoints) as a custom-geometry shape."""
    xs = [int(round(p[0])) for p in points]
    ys = [int(round(p[1])) for p in points]
    offx, offy = min(xs), min(ys)
    cx = max(1, max(xs) - offx)
    cy = max(1, max(ys) - offy)
    moveto = '<a:moveTo><a:pt x="%d" y="%d"/></a:moveTo>' % (xs[0] - offx, ys[0] - offy)
    segs = "".join('<a:lnTo><a:pt x="%d" y="%d"/></a:lnTo>' % (x - offx, y - offy)
                  for x, y in zip(xs[1:], ys[1:]))
    ln = build_ln_xml(style)
    cid = next_id()
    xml = (
        '<p:sp %s>'
        '<p:nvSpPr><p:cNvPr id="%d" name="edge"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr>'
        '<a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
        '<a:custGeom><a:avLst/><a:gdLst/><a:rectLst/><a:pathLst>'
        '<a:path>%s%s</a:path></a:pathLst></a:custGeom>'
        '<a:noFill/>%s'
        '</p:spPr><p:style/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        '</p:sp>'
    ) % (nsdecls("a", "p"), cid, offx, offy, cx, cy, moveto, segs, ln)
    sp = parse_xml(xml)
    slide.shapes._spTree.append(sp)
    return sp


def set_vertical_text(tf):
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("vert", "vert270")


# ---------------------------------------------------------------------------
# Geometry: resolve absolute coordinates for nested containers
# ---------------------------------------------------------------------------

def parse_float(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


class Cell:
    __slots__ = ("id", "parent", "style", "value", "is_edge", "is_vertex",
                 "x", "y", "w", "h", "source", "target", "points",
                 "src_point", "tgt_point", "ax", "ay")

    def __init__(self, el):
        self.id = el.get("id")
        self.parent = el.get("parent", "1")
        self.style = parse_style(el.get("style", ""))
        self.value = el.get("value", "")
        self.is_edge = el.get("edge") == "1"
        self.is_vertex = el.get("vertex") == "1"
        self.x = self.y = self.w = self.h = 0.0
        self.source = el.get("source")
        self.target = el.get("target")
        self.points = []
        self.src_point = None
        self.tgt_point = None
        geom = el.find("mxGeometry")
        if geom is not None:
            self.x = parse_float(geom.get("x"))
            self.y = parse_float(geom.get("y"))
            self.w = parse_float(geom.get("width"), 0.0)
            self.h = parse_float(geom.get("height"), 0.0)
            if self.is_edge:
                src = geom.find("mxPoint[@as='sourcePoint']")
                tgt = geom.find("mxPoint[@as='targetPoint']")
                if src is not None:
                    self.src_point = (parse_float(src.get("x")), parse_float(src.get("y")))
                if tgt is not None:
                    self.tgt_point = (parse_float(tgt.get("x")), parse_float(tgt.get("y")))
                for p in geom.iterfind(".//mxPoint[@as='array']") or []:
                    pass
                arr = geom.find("Array[@as='points']")
                if arr is not None:
                    for p in arr.findall("mxPoint"):
                        self.points.append((parse_float(p.get("x")), parse_float(p.get("y"))))


def resolve_geometry(cells):
    """Compute absolute x,y,w,h for every vertex by summing parent offsets."""
    by_id = {c.id: c for c in cells}

    def abs_of(cell, seen=None):
        seen = seen or set()
        if cell.id in seen:
            return 0.0, 0.0
        seen.add(cell.id)
        if cell.parent and cell.parent in by_id:
            parent = by_id[cell.parent]
            if parent.is_vertex:
                px, py = abs_of(parent, seen)
                return cell.x + px, cell.y + py
        return cell.x, cell.y

    for c in cells:
        if c.is_vertex:
            ax, ay = abs_of(c)
            c.ax, c.ay = ax, ay  # absolute origin
        else:
            c.ax, c.ay = c.x, c.y
    return by_id


# ---------------------------------------------------------------------------
# Edge endpoint clipping
# ---------------------------------------------------------------------------

def rect_center(ax, ay, w, h):
    return (ax + w / 2.0, ay + h / 2.0)


def clip_to_rect(ax, ay, w, h, tx, ty):
    """Point on the boundary of rect (ax,ay,w,h) along the line to (tx,ty)."""
    cx, cy = rect_center(ax, ay, w, h)
    dx, dy = tx - cx, ty - cy
    if dx == 0 and dy == 0:
        return cx, cy
    sx = (w / 2.0) / abs(dx) if dx != 0 else float("inf")
    sy = (h / 2.0) / abs(dy) if dy != 0 else float("inf")
    s = min(sx, sy)
    return cx + dx * s, cy + dy * s


# ---------------------------------------------------------------------------
# Building the slide
# ---------------------------------------------------------------------------

def build_slide(prs, model, args):
    cells = [Cell(c) for c in model.findall(".//mxCell")]
    by_id = resolve_geometry(cells)
    vertices = [c for c in cells if c.is_vertex and c.w > 0 and c.h > 0]
    edges = [c for c in cells if c.is_edge]

    if not vertices:
        blank = prs.slides.add_slide(prs.slide_layouts[6])
        return

    # bounding box of all vertices
    min_x = min(c.ax for c in vertices)
    min_y = min(c.ay for c in vertices)
    max_x = max(c.ax + c.w for c in vertices)
    max_y = max(c.ay + c.h for c in vertices)
    bw, bh = max_x - min_x, max_y - min_y

    # slide size + fit
    sw, sh = prs.slide_width, prs.slide_height
    margin = Emu(int(args.margin * EMU_PER_IN))
    avail_w = sw - 2 * margin
    avail_h = sh - 2 * margin
    scale = min(avail_w / max(bw, 1), avail_h / max(bh, 1))
    cpad_x = (avail_w - bw * scale) / 2.0
    cpad_y = (avail_h - bh * scale) / 2.0

    def mx(px):
        return int(round((px - min_x) * scale)) + margin + int(cpad_x)

    def my(py):
        return int(round((py - min_y) * scale)) + margin + int(cpad_y)

    def mw(pw):
        return int(round(pw * scale))

    def mh(ph):
        return int(round(ph * scale))

    pt_per_px = scale * 72.0 / EMU_PER_IN

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background (from model attribute if present)
    bg = model.get("background")
    if bg and bg != "none":
        rgb = hexcolor(bg)
        if rgb is not None:
            slide.background.fill.solid()
            slide.background.fore_color.rgb = rgb

    shape_map = {}  # cell id -> pptx shape (for connector attach by id, optional)

    # ---- vertices ----
    for c in vertices:
        left, top = mx(c.ax), my(c.ay)
        w, h = mw(c.w), mh(c.h)
        kind = shape_kind(c.style)
        if kind is None:
            shape = slide.shapes.add_textbox(left, top, max(w, Emu(36000)), max(h, Emu(18000)))
        else:
            shape = slide.shapes.add_shape(kind, left, top, max(w, Emu(18000)), max(h, Emu(18000)))
        shape_map[c.id] = shape
        apply_shape_style(shape, c.style)
        apply_text(shape, c.value, c.style, pt_per_px, args.font_scale,
                   args.min_font, args.max_font)

    # ---- edges ----
    for e in edges:
        pts = compute_edge_points(e, by_id, mx, my)
        if len(pts) < 2:
            continue
        if len(pts) == 2:
            (x1, y1), (x2, y2) = pts
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            set_connector_arrows(conn, e.style)
            # width
            w = e.style.get("strokeWidth")
            if w:
                try:
                    conn.line.width = Pt(float(w))
                except ValueError:
                    pass
        else:
            add_polyline(slide, pts, e.style)
        # edge label
        if e.value and e.value.strip():
            lx = sum(p[0] for p in pts) / len(pts)
            ly = sum(p[1] for p in pts) / len(pts)
            add_edge_label(slide, e.value, e.style, lx, ly, pt_per_px, args.font_scale,
                           args.min_font, args.max_font)


def compute_edge_points(e, by_id, mx, my):
    """Return list of (emu_x, emu_y) points for an edge (begin, waypoints, end)."""
    src = by_id.get(e.source) if e.source else None
    tgt = by_id.get(e.target) if e.target else None
    wps = [(mx(px), my(py)) for (px, py) in e.points]

    if src and src.is_vertex and tgt and tgt.is_vertex:
        sx, sy, sw, sh = src.ax, src.ay, src.w, src.h
        tx, ty, tw, th = tgt.ax, tgt.ay, tgt.w, tgt.h
        if wps:
            start = clip_to_rect(sx, sy, sw, sh, e.points[0][0], e.points[0][1])
            end = clip_to_rect(tx, ty, tw, th, e.points[-1][0], e.points[-1][1])
        else:
            sc = rect_center(sx, sy, sw, sh)
            tc = rect_center(tx, ty, tw, th)
            start = clip_to_rect(sx, sy, sw, sh, tc[0], tc[1])
            end = clip_to_rect(tx, ty, tw, th, sc[0], sc[1])
        return [(mx(start[0]), my(start[1]))] + wps + [(mx(end[0]), my(end[1]))]

    # fall back to explicit points
    out = []
    if e.src_point:
        out.append((mx(e.src_point[0]), my(e.src_point[1])))
    out.extend(wps)
    if e.tgt_point:
        out.append((mx(e.tgt_point[0]), my(e.tgt_point[1])))
    return out


def apply_shape_style(shape, style):
    fc = style.get("fillColor")
    if fc == "none":
        shape.fill.background()
    else:
        rgb = hexcolor(fc, RGBColor(0xFF, 0xFF, 0xFF))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    sc = style.get("strokeColor")
    if sc == "none":
        shape.line.fill.background()
    else:
        rgb = hexcolor(sc, RGBColor(0, 0, 0))
        shape.line.color.rgb = rgb
        w = style.get("strokeWidth")
        if w:
            try:
                shape.line.width = Pt(float(w))
            except ValueError:
                pass
    # disable autosize shadow default
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def apply_text(shape, value, style, pt_per_px, font_scale, min_pt, max_pt):
    lines = html_to_lines(value)
    if not lines or (len(lines) == 1 and not lines[0].strip()):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(2))
    align = style.get("align", "center")
    palign = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    fcolor = hexcolor(style.get("fontColor"), RGBColor(0x1F, 0x1F, 0x1F))
    fstyle = 0
    try:
        fstyle = int(style.get("fontStyle", "0"))
    except ValueError:
        fstyle = 0
    bold = bool(fstyle & 1)
    italic = bool(fstyle & 2)
    underline = bool(fstyle & 4)
    try:
        fpx = float(style.get("fontSize", "12"))
    except ValueError:
        fpx = 12.0
    pt = max(min_pt, min(max_pt, fpx * pt_per_px * font_scale))
    fname = style.get("fontFamily")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = palign
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(pt)
        f.bold = bold
        f.italic = italic
        f.underline = underline
        if fcolor:
            f.color.rgb = fcolor
        if fname:
            f.name = fname
    va = style.get("verticalAlign", "middle")
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM}.get(va, MSO_ANCHOR.MIDDLE)
    if style.get("horizontal", "1") == "0":
        set_vertical_text(tf)


def add_edge_label(slide, value, style, cx, cy, pt_per_px, font_scale, min_pt, max_pt):
    text = " ".join(html_to_lines(value))
    text = text.strip()
    if not text:
        return
    box_w = Emu(min(max(len(text) * int(9 * pt_per_px * font_scale * 1.6), 200000), 1500000))
    box_h = Emu(max(int(16 * pt_per_px * font_scale * 1.6), 180000))
    tb = slide.shapes.add_textbox(int(cx - box_w / 2), int(cy - box_h / 2), box_w, box_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    f = run.font
    try:
        fpx = float(style.get("fontSize", "11"))
    except ValueError:
        fpx = 11.0
    pt = max(min_pt, min(max_pt, fpx * pt_per_px * font_scale))
    f.size = Pt(pt)
    f.bold = True
    fcolor = hexcolor(style.get("fontColor"))
    if fcolor:
        f.color.rgb = fcolor
    bg = style.get("labelBackgroundColor")
    if bg and bg != "none":
        rgb = hexcolor(bg)
        if rgb:
            tb.fill.solid()
            tb.fill.fore_color.rgb = rgb
            tb.line.fill.background()
    else:
        tb.fill.background()
        tb.line.fill.background()


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

SIZE_PRESETS = {
    "16:9": (Emu(int(13.333 * EMU_PER_IN)), Emu(int(7.5 * EMU_PER_IN))),
    "4:3": (Emu(int(10 * EMU_PER_IN)), Emu(int(7.5 * EMU_PER_IN))),
    "16:10": (Emu(int(13.333 * EMU_PER_IN)), Emu(int(8.33 * EMU_PER_IN))),
    "A4-landscape": (Emu(int(11.69 * EMU_PER_IN)), Emu(int(8.27 * EMU_PER_IN))),
}


def main():
    ap = argparse.ArgumentParser(description="Convert .drawio to a native editable .pptx")
    ap.add_argument("input", help="input .drawio file")
    ap.add_argument("-o", "--output", required=True, help="output .pptx file")
    ap.add_argument("--size", default="16:9", choices=list(SIZE_PRESETS),
                    help="slide size preset (default 16:9)")
    ap.add_argument("--landscape", action="store_true",
                    help="auto-pick 16:9 (landscape) if the diagram is wider than tall, else 4:3")
    ap.add_argument("--margin", type=float, default=0.3, help="slide margin in inches (default 0.3)")
    ap.add_argument("--font-scale", type=float, default=1.0, help="text size multiplier (default 1.0)")
    ap.add_argument("--min-font", type=float, default=7.0, help="minimum font pt (default 7)")
    ap.add_argument("--max-font", type=float, default=40.0, help="maximum font pt (default 40)")
    args = ap.parse_args()

    diagrams = load_diagrams(args.input)
    prs = Presentation()

    if args.landscape:
        # peek at first diagram aspect
        cells = [Cell(c) for c in diagrams[0][1].findall(".//mxCell")]
        resolve_geometry(cells)
        vs = [c for c in cells if c.is_vertex and c.w > 0 and c.h > 0]
        if vs:
            w = max(c.ax + c.w for c in vs) - min(c.ax for c in vs)
            h = max(c.ay + c.h for c in vs) - min(c.ay for c in vs)
            args.size = "16:9" if w >= h else "4:3"

    prs.slide_width, prs.slide_height = SIZE_PRESETS[args.size]

    for name, model in diagrams:
        build_slide(prs, model, args)

    prs.save(args.output)
    print("Wrote %s (%d slide(s))" % (args.output, len(diagrams)))


if __name__ == "__main__":
    main()
